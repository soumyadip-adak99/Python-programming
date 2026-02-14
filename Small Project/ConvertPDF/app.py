import os
import uuid
import threading
import time
import atexit
import zipfile
import io
from flask import Flask, render_template, request, jsonify, url_for, make_response, send_file
from werkzeug.utils import secure_filename
import img2pdf
from PIL import Image
from docx import Document
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from PyPDF2 import PdfReader, PdfWriter, PdfMerger

# Optional imports (may not be installed)
try:
    from pdf2image import convert_from_path
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

app = Flask(__name__)
app.config["SECRET_KEY"] = "your-secret-key-here"
app.config["UPLOAD_FOLDER"] = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "uploads"
)
app.config["OUTPUT_FOLDER"] = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "outputs"
)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100MB max file size

# Allowed extensions
ALLOWED_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "bmp", "txt", "docx", "html", "webp", "pdf", "xlsx", "xls", "pptx", "ppt"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "bmp", "webp"}

# Ensure directories exist
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
os.makedirs(app.config["OUTPUT_FOLDER"], exist_ok=True)


# ============================================================
# Scheduler to clean up files every 7 minutes
# ============================================================
class FileCleanupScheduler:
    def __init__(self, folders, interval_minutes=7):
        self.folders = folders
        self.interval = interval_minutes * 60
        self.running = False
        self.thread = None

    def cleanup_files(self):
        """Delete all files in the folders"""
        try:
            for folder in self.folders:
                if os.path.exists(folder):
                    for filename in os.listdir(folder):
                        file_path = os.path.join(folder, filename)
                        try:
                            if os.path.isfile(file_path):
                                os.remove(file_path)
                                print(f"[Scheduler] Deleted: {filename}")
                        except Exception as e:
                            print(f"[Scheduler] Error deleting {filename}: {e}")
            print(f"[Scheduler] Cleanup completed at {time.strftime('%H:%M:%S')}")
        except Exception as e:
            print(f"[Scheduler] Cleanup error: {e}")

    def run(self):
        while self.running:
            time.sleep(self.interval)
            if self.running:
                self.cleanup_files()

    def start(self):
        if not self.running:
            self.running = True
            self.thread = threading.Thread(target=self.run, daemon=True)
            self.thread.start()
            print(f"[Scheduler] Started - will clean up files every 7 minutes")

    def stop(self):
        self.running = False
        print("[Scheduler] Stopped")


# Create and start the scheduler
file_scheduler = FileCleanupScheduler(
    [app.config["UPLOAD_FOLDER"], app.config["OUTPUT_FOLDER"]], 
    interval_minutes=7
)


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_extension(filename):
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


def cleanup_folder(folder):
    """Delete all files in a folder"""
    try:
        if os.path.exists(folder):
            for filename in os.listdir(folder):
                file_path = os.path.join(folder, filename)
                try:
                    if os.path.isfile(file_path):
                        os.remove(file_path)
                except:
                    pass
    except:
        pass


# ============================================================
# CONVERT TO PDF FUNCTIONS
# ============================================================

def convert_image_to_pdf(image_path, output_path):
    """Convert image files to PDF - preserves original size and quality"""
    try:
        temp_path = None
        convert_path = image_path

        with Image.open(image_path) as img:
            if img.mode in ("RGBA", "LA", "P"):
                rgb_img = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                if img.mode in ("RGBA", "LA"):
                    rgb_img.paste(img, mask=img.split()[-1])
                else:
                    rgb_img.paste(img)
                temp_path = image_path + "_temp.jpg"
                rgb_img.save(temp_path, "JPEG", quality=100)
                convert_path = temp_path

        with open(output_path, "wb") as f:
            f.write(img2pdf.convert(convert_path))

        if temp_path and os.path.exists(temp_path):
            os.remove(temp_path)

        return True
    except Exception as e:
        print(f"Error converting image: {e}")
        return False


def convert_text_to_pdf(text_path, output_path):
    """Convert text files to PDF"""
    try:
        with open(text_path, "r", encoding="utf-8", errors="ignore") as f:
            text_content = f.read()

        doc = SimpleDocTemplate(
            output_path, pagesize=A4,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
        )
        styles = getSampleStyleSheet()
        story = []

        paragraphs = text_content.split("\n")
        for para in paragraphs:
            if para.strip():
                story.append(Paragraph(para, styles["Normal"]))
            story.append(Spacer(1, 12))

        if not story:
            story.append(Paragraph("Empty document", styles["Normal"]))

        doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting text: {e}")
        return False


def convert_docx_to_pdf(docx_path, output_path):
    """Convert DOCX files to PDF"""
    try:
        doc = Document(docx_path)
        pdf_doc = SimpleDocTemplate(
            output_path, pagesize=A4,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
        )
        styles = getSampleStyleSheet()
        story = []

        for para in doc.paragraphs:
            if para.text.strip():
                story.append(Paragraph(para.text, styles["Normal"]))
            story.append(Spacer(1, 12))

        if not story:
            story.append(Paragraph("Empty document", styles["Normal"]))

        pdf_doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting DOCX: {e}")
        return False


def convert_html_to_pdf(html_path, output_path):
    """Convert HTML files to PDF"""
    try:
        import re

        with open(html_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        text_content = re.sub(r"<[^>]+>", "", html_content)
        text_content = re.sub(r"\s+", " ", text_content).strip()

        doc = SimpleDocTemplate(
            output_path, pagesize=A4,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
        )
        styles = getSampleStyleSheet()
        story = []

        words = text_content.split()
        chunks = [" ".join(words[i : i + 100]) for i in range(0, len(words), 100)]

        for chunk in chunks:
            if chunk.strip():
                story.append(Paragraph(chunk, styles["Normal"]))
                story.append(Spacer(1, 12))

        if not story:
            story.append(Paragraph("Empty document", styles["Normal"]))

        doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting HTML: {e}")
        return False


def convert_excel_to_pdf(excel_path, output_path):
    """Convert Excel files to PDF"""
    if not OPENPYXL_AVAILABLE:
        return False
    try:
        wb = load_workbook(excel_path)
        sheet = wb.active

        doc = SimpleDocTemplate(
            output_path, pagesize=A4,
            rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36,
        )
        styles = getSampleStyleSheet()
        story = []

        # Extract data from sheet
        data = []
        for row in sheet.iter_rows(values_only=True):
            data.append([str(cell) if cell is not None else "" for cell in row])

        if data:
            table = Table(data)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
            ]))
            story.append(table)
        else:
            story.append(Paragraph("Empty spreadsheet", styles["Normal"]))

        doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting Excel: {e}")
        return False


def convert_pptx_to_pdf(pptx_path, output_path):
    """Convert PowerPoint files to PDF"""
    if not PPTX_AVAILABLE:
        return False
    try:
        prs = Presentation(pptx_path)
        
        doc = SimpleDocTemplate(
            output_path, pagesize=A4,
            rightMargin=72, leftMargin=72, topMargin=72, bottomMargin=72,
        )
        styles = getSampleStyleSheet()
        story = []

        for i, slide in enumerate(prs.slides, 1):
            story.append(Paragraph(f"<b>Slide {i}</b>", styles["Heading2"]))
            story.append(Spacer(1, 12))
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    story.append(Paragraph(shape.text, styles["Normal"]))
                    story.append(Spacer(1, 6))
            
            story.append(Spacer(1, 24))

        if not story:
            story.append(Paragraph("Empty presentation", styles["Normal"]))

        doc.build(story)
        return True
    except Exception as e:
        print(f"Error converting PowerPoint: {e}")
        return False


def convert_to_pdf(input_path, output_path, file_extension):
    """Main conversion function"""
    converters = {
        "png": convert_image_to_pdf,
        "jpg": convert_image_to_pdf,
        "jpeg": convert_image_to_pdf,
        "gif": convert_image_to_pdf,
        "bmp": convert_image_to_pdf,
        "webp": convert_image_to_pdf,
        "txt": convert_text_to_pdf,
        "docx": convert_docx_to_pdf,
        "html": convert_html_to_pdf,
        "xlsx": convert_excel_to_pdf,
        "xls": convert_excel_to_pdf,
        "pptx": convert_pptx_to_pdf,
        "ppt": convert_pptx_to_pdf,
    }

    converter = converters.get(file_extension)
    if converter:
        return converter(input_path, output_path)
    return False


# ============================================================
# PDF MANIPULATION FUNCTIONS
# ============================================================

def merge_pdfs(pdf_paths, output_path):
    """Merge multiple PDF files into one"""
    try:
        merger = PdfMerger()
        for pdf_path in pdf_paths:
            merger.append(pdf_path)
        merger.write(output_path)
        merger.close()
        return True
    except Exception as e:
        print(f"Error merging PDFs: {e}")
        return False


def split_pdf(pdf_path, output_path, pages):
    """Split PDF - extract specific pages"""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        
        total_pages = len(reader.pages)
        
        # Parse pages string like "1,3,5-7"
        page_numbers = set()
        for part in pages.split(","):
            part = part.strip()
            if "-" in part:
                start, end = part.split("-")
                for p in range(int(start), min(int(end) + 1, total_pages + 1)):
                    page_numbers.add(p)
            else:
                page_numbers.add(int(part))
        
        for page_num in sorted(page_numbers):
            if 1 <= page_num <= total_pages:
                writer.add_page(reader.pages[page_num - 1])
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        return True
    except Exception as e:
        print(f"Error splitting PDF: {e}")
        return False


def compress_pdf(pdf_path, output_path, quality="medium"):
    """Compress PDF to reduce file size"""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
        
        # Set compression level
        if quality == "high":
            # Minimal compression, best quality
            pass
        elif quality == "low":
            # Maximum compression
            for page in writer.pages:
                page.compress_content_streams()
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        return True
    except Exception as e:
        print(f"Error compressing PDF: {e}")
        return False


def rotate_pdf(pdf_path, output_path, angle):
    """Rotate all pages in PDF"""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.rotate(int(angle))
            writer.add_page(page)
        
        with open(output_path, "wb") as f:
            writer.write(f)
        
        return True
    except Exception as e:
        print(f"Error rotating PDF: {e}")
        return False


def pdf_to_images(pdf_path, output_folder, format="jpg"):
    """Convert PDF pages to images"""
    if not PDF2IMAGE_AVAILABLE:
        return None
    try:
        images = convert_from_path(pdf_path)
        image_paths = []
        
        for i, image in enumerate(images, 1):
            image_filename = f"page_{i}.{format}"
            image_path = os.path.join(output_folder, image_filename)
            image.save(image_path, format.upper() if format != "jpg" else "JPEG")
            image_paths.append(image_path)
        
        return image_paths
    except Exception as e:
        print(f"Error converting PDF to images: {e}")
        return None


# ============================================================
# ROUTES
# ============================================================

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/tools")
def get_tools():
    """Return available tools"""
    tools = [
        {"id": "convert", "name": "Convert to PDF", "description": "Images, DOCX, TXT, HTML, Excel, PPT"},
        {"id": "merge", "name": "Merge PDFs", "description": "Combine multiple PDFs into one"},
        {"id": "split", "name": "Split PDF", "description": "Extract specific pages"},
        {"id": "compress", "name": "Compress PDF", "description": "Reduce file size"},
        {"id": "rotate", "name": "Rotate PDF", "description": "Rotate pages 90°, 180°, 270°"},
    ]
    if PDF2IMAGE_AVAILABLE:
        tools.insert(1, {"id": "pdf-to-image", "name": "PDF to Image", "description": "Convert PDF pages to JPG/PNG"})
    return jsonify({"tools": tools})


@app.route("/upload", methods=["POST"])
@app.route("/api/convert", methods=["POST"])
def upload_file():
    """Convert file to PDF"""
    if "file" not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400

    if file and allowed_file(file.filename):
        cleanup_folder(app.config["OUTPUT_FOLDER"])

        original_filename = secure_filename(file.filename)
        file_extension = get_file_extension(original_filename)
        unique_id = str(uuid.uuid4())[:8]

        upload_path = os.path.join(
            app.config["UPLOAD_FOLDER"], f"{unique_id}_{original_filename}"
        )
        file.save(upload_path)

        pdf_filename = f"{unique_id}_{original_filename.rsplit('.', 1)[0]}.pdf"
        output_path = os.path.join(app.config["OUTPUT_FOLDER"], pdf_filename)

        success = convert_to_pdf(upload_path, output_path, file_extension)

        try:
            os.remove(upload_path)
        except:
            pass

        if success:
            return jsonify({
                "success": True,
                "message": "File converted successfully!",
                "filename": pdf_filename,
                "download_url": url_for("download_file", filename=pdf_filename),
            })
        else:
            return jsonify({"error": "Failed to convert file"}), 500

    return jsonify({"error": "File type not allowed"}), 400


@app.route("/api/merge", methods=["POST"])
def merge_files():
    """Merge multiple PDF files"""
    if "files" not in request.files:
        return jsonify({"error": "No files provided"}), 400
    
    files = request.files.getlist("files")
    if len(files) < 2:
        return jsonify({"error": "Need at least 2 PDF files to merge"}), 400
    
    cleanup_folder(app.config["OUTPUT_FOLDER"])
    unique_id = str(uuid.uuid4())[:8]
    pdf_paths = []
    
    try:
        for file in files:
            if file and file.filename.lower().endswith('.pdf'):
                filename = secure_filename(file.filename)
                save_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}")
                file.save(save_path)
                pdf_paths.append(save_path)
        
        if len(pdf_paths) < 2:
            return jsonify({"error": "Need at least 2 valid PDF files"}), 400
        
        output_filename = f"{unique_id}_merged.pdf"
        output_path = os.path.join(app.config["OUTPUT_FOLDER"], output_filename)
        
        success = merge_pdfs(pdf_paths, output_path)
        
        # Cleanup uploaded files
        for path in pdf_paths:
            try:
                os.remove(path)
            except:
                pass
        
        if success:
            return jsonify({
                "success": True,
                "message": "PDFs merged successfully!",
                "filename": output_filename,
                "download_url": url_for("download_file", filename=output_filename),
            })
        else:
            return jsonify({"error": "Failed to merge PDFs"}), 500
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/split", methods=["POST"])
def split_file():
    """Split PDF by pages"""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    pages = request.form.get("pages", "1")
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({"error": "Only PDF files can be split"}), 400
    
    cleanup_folder(app.config["OUTPUT_FOLDER"])
    unique_id = str(uuid.uuid4())[:8]
    
    try:
        filename = secure_filename(file.filename)
        upload_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}")
        file.save(upload_path)
        
        output_filename = f"{unique_id}_split.pdf"
        output_path = os.path.join(app.config["OUTPUT_FOLDER"], output_filename)
        
        success = split_pdf(upload_path, output_path, pages)
        
        try:
            os.remove(upload_path)
        except:
            pass
        
        if success:
            return jsonify({
                "success": True,
                "message": "PDF split successfully!",
                "filename": output_filename,
                "download_url": url_for("download_file", filename=output_filename),
            })
        else:
            return jsonify({"error": "Failed to split PDF"}), 500
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/compress", methods=["POST"])
def compress_file():
    """Compress PDF"""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    quality = request.form.get("quality", "medium")
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({"error": "Only PDF files can be compressed"}), 400
    
    cleanup_folder(app.config["OUTPUT_FOLDER"])
    unique_id = str(uuid.uuid4())[:8]
    
    try:
        filename = secure_filename(file.filename)
        upload_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}")
        file.save(upload_path)
        
        original_size = os.path.getsize(upload_path)
        
        output_filename = f"{unique_id}_compressed.pdf"
        output_path = os.path.join(app.config["OUTPUT_FOLDER"], output_filename)
        
        success = compress_pdf(upload_path, output_path, quality)
        
        try:
            os.remove(upload_path)
        except:
            pass
        
        if success:
            new_size = os.path.getsize(output_path)
            reduction = round((1 - new_size / original_size) * 100, 1)
            
            return jsonify({
                "success": True,
                "message": f"PDF compressed! Reduced by {reduction}%",
                "filename": output_filename,
                "download_url": url_for("download_file", filename=output_filename),
                "original_size": original_size,
                "compressed_size": new_size,
                "reduction_percent": reduction,
            })
        else:
            return jsonify({"error": "Failed to compress PDF"}), 500
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/rotate", methods=["POST"])
def rotate_file():
    """Rotate PDF pages"""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    angle = request.form.get("angle", "90")
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({"error": "Only PDF files can be rotated"}), 400
    
    cleanup_folder(app.config["OUTPUT_FOLDER"])
    unique_id = str(uuid.uuid4())[:8]
    
    try:
        filename = secure_filename(file.filename)
        upload_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}")
        file.save(upload_path)
        
        output_filename = f"{unique_id}_rotated.pdf"
        output_path = os.path.join(app.config["OUTPUT_FOLDER"], output_filename)
        
        success = rotate_pdf(upload_path, output_path, angle)
        
        try:
            os.remove(upload_path)
        except:
            pass
        
        if success:
            return jsonify({
                "success": True,
                "message": f"PDF rotated by {angle}°!",
                "filename": output_filename,
                "download_url": url_for("download_file", filename=output_filename),
            })
        else:
            return jsonify({"error": "Failed to rotate PDF"}), 500
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/pdf-to-image", methods=["POST"])
def convert_pdf_to_images():
    """Convert PDF to images"""
    if not PDF2IMAGE_AVAILABLE:
        return jsonify({"error": "PDF to image conversion not available. Install poppler."}), 400
    
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    format = request.form.get("format", "jpg")
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({"error": "Only PDF files can be converted to images"}), 400
    
    cleanup_folder(app.config["OUTPUT_FOLDER"])
    unique_id = str(uuid.uuid4())[:8]
    
    try:
        filename = secure_filename(file.filename)
        upload_path = os.path.join(app.config["UPLOAD_FOLDER"], f"{unique_id}_{filename}")
        file.save(upload_path)
        
        image_paths = pdf_to_images(upload_path, app.config["OUTPUT_FOLDER"], format)
        
        try:
            os.remove(upload_path)
        except:
            pass
        
        if image_paths:
            # Create ZIP file
            zip_filename = f"{unique_id}_images.zip"
            zip_path = os.path.join(app.config["OUTPUT_FOLDER"], zip_filename)
            
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for img_path in image_paths:
                    zipf.write(img_path, os.path.basename(img_path))
                    os.remove(img_path)
            
            return jsonify({
                "success": True,
                "message": f"PDF converted to {len(image_paths)} images!",
                "filename": zip_filename,
                "download_url": url_for("download_file", filename=zip_filename),
                "page_count": len(image_paths),
            })
        else:
            return jsonify({"error": "Failed to convert PDF to images"}), 500
            
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/pdf-info", methods=["POST"])
def get_pdf_info():
    """Get PDF information (page count, etc.)"""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    
    if not file.filename.lower().endswith('.pdf'):
        return jsonify({"error": "Only PDF files accepted"}), 400
    
    try:
        reader = PdfReader(file)
        return jsonify({
            "success": True,
            "page_count": len(reader.pages),
            "filename": file.filename,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/download/<filename>")
def download_file(filename):
    """Download and delete file after sending"""
    safe_filename = secure_filename(filename)
    file_path = os.path.join(app.config["OUTPUT_FOLDER"], safe_filename)

    if not os.path.exists(file_path):
        return jsonify({"error": "File not found"}), 404

    try:
        with open(file_path, "rb") as f:
            file_data = f.read()

        os.remove(file_path)

        response = make_response(file_data)
        
        if safe_filename.endswith('.pdf'):
            response.headers["Content-Type"] = "application/pdf"
        elif safe_filename.endswith('.zip'):
            response.headers["Content-Type"] = "application/zip"
        else:
            response.headers["Content-Type"] = "application/octet-stream"
            
        response.headers["Content-Disposition"] = f'attachment; filename="{safe_filename}"'
        response.headers["Content-Length"] = len(file_data)
        response.headers["Cache-Control"] = "no-cache, no-store, must-revalidate"
        response.headers["X-Content-Type-Options"] = "nosniff"

        return response
    except Exception as e:
        print(f"Error downloading file: {e}")
        return jsonify({"error": "Download failed"}), 500


@app.route("/health")
def health():
    return jsonify({
        "status": "healthy",
        "features": {
            "pdf2image": PDF2IMAGE_AVAILABLE,
            "excel": OPENPYXL_AVAILABLE,
            "powerpoint": PPTX_AVAILABLE,
        }
    })


# Register cleanup on app shutdown
atexit.register(file_scheduler.stop)


if __name__ == "__main__":
    file_scheduler.start()
    app.run(host="0.0.0.0", port=5001, debug=True, use_reloader=False)
